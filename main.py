import logging
import traceback
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    ContextTypes,
    CallbackQueryHandler,
    MessageHandler,
    filters,
    ConversationHandler,
)
from telegram.error import TimedOut
from sheets import read_sheet, get_sheet_title
from pdf_generator import generate_pdf
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os
import subprocess
import requests
from PIL import Image as PIL_Image
from io import BytesIO
import uuid
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Настройка логирования на запись в файл и консоль
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler("bot_log.txt"),
        logging.StreamHandler()
    ]
)

# Проверка наличия необходимых папок
ORIGINAL_TEMPLATES_DIR = 'original_templates'
UPDATED_TEMPLATES_DIR = 'updated_templates'

os.makedirs(ORIGINAL_TEMPLATES_DIR, exist_ok=True)
os.makedirs(UPDATED_TEMPLATES_DIR, exist_ok=True)

logging.info("Настройка логирования завершена. Логирование начато.")

def get_template_names(directory=ORIGINAL_TEMPLATES_DIR):
    """
    Получает список шаблонов из папки original_templates.
    """
    template_files = os.listdir(directory)
    template_names = [os.path.splitext(file)[0] for file in template_files if file.endswith('.pptx')]
    return template_names

def convert_drive_url(url):
    """
    Преобразует ссылку Google Drive в прямую ссылку для скачивания.
    """
    if "drive.google.com" in url:
        try:
            # Извлекаем FILE_ID из ссылки
            file_id = url.split("/d/")[1].split("/")[0]
            direct_url = f"https://drive.google.com/uc?export=download&id={file_id}"
            return direct_url
        except IndexError:
            logging.error(f"Невозможно извлечь ID файла из URL: {url}")
            return url
    return url

def download_photo(url):
    try:
        direct_url = convert_drive_url(url)
        response = requests.get(direct_url)
        response.raise_for_status()
        try:
            image = PIL_Image.open(BytesIO(response.content))
            image.load()
            unique_filename = f"/tmp/photo_{uuid.uuid4().hex}.jpg"
            image.save(unique_filename)
            return unique_filename
        except (IOError, PIL_Image.UnidentifiedImageError) as img_err:
            logging.error(f"Невалидный файл изображения: {img_err}")
            return None
    except requests.exceptions.RequestException as e:
        logging.error(f"Ошибка при загрузке фото: {e}")
        return None


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    logging.info('Команда /start вызвана')
    keyboard = [
        [InlineKeyboardButton("📋 Получить", callback_data='get_contractors')],
        [InlineKeyboardButton("➕ Загрузить шаблон", callback_data='upload_template')],
        [InlineKeyboardButton("📑 Показать шаблоны", callback_data='show_templates')],
        [InlineKeyboardButton("❌ Удалить шаблон", callback_data='delete_template')],
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.message.reply_text('Добро пожаловать! Выберите действие:', reply_markup=reply_markup)

async def get_contractors(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    try:
        logging.info('Команда /get_contractors вызвана')
        spreadsheet_id = '1syI_hA_2n-4dygYYHZ5vt0_l2ys2uZXDlD3nONglWOU'
        range_name = 'Sheet1!A2:J'  
        values = read_sheet(spreadsheet_id, range_name)
        if not values:
            await update.reply_text('Нет данных в таблице.')
        else:
            logging.info('Таблица данных:\n%s', '\n'.join(['\t'.join(row) for row in values]))
    
            table_header = "Имя            | Фамилия | Город   | Стоимость | Часы | Мин. часы | Трансфер | Instagram         | Портфолио          | VK\n"
            table_divider = "---------------|---------|---------|-----------|------|-----------|----------|-------------------|--------------------|----\n"
            table_rows = ""
            for row in values:
                name = row[1] if len(row) > 1 and row[1] else 'Не указано'
                surname = row[1].split()[1] if len(row[1].split()) > 1 else 'N/A'
                city = row[2] if len(row) > 2 else 'N/A'
                cost = row[3] if len(row) > 3 else 'N/A'
                hours = row[4] if len(row) > 4 else 'N/A'
                min_hours = row[5] if len(row) > 5 else 'N/A'
                transfer = row[6] if len(row) > 6 else 'N/A'
                instagram = extract_username(row[7]) if len(row) > 7 else 'N/A'
                portfolio = row[8] if len(row) > 8 else 'N/A'
                vk = extract_username(row[9]) if len(row) > 9 else 'N/A'
                table_rows += f"{name:<15} | {surname:<7} | {city:<7} | {cost:<9} | {hours:<4} | {min_hours:<9} | {transfer:<8} | {instagram:<17} | {portfolio:<18} | {vk}\n"
            
            table = f"```\n{table_header}{table_divider}{table_rows}```"
            await update.reply_text(table, parse_mode='Markdown')
    except Exception as e:
        logging.error(f'Ошибка при получении подрядчиков: {e}')
        logging.error(traceback.format_exc())
        await update.reply_text('Произошла ошибка при обработке запроса.')

async def show_personnel_list(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    personnel_data = context.user_data.get('personnel_data', [])
    selected_personnel = context.user_data.get('selected_personnel', set())
    keyboard = []
    for row in personnel_data:
        button_text = (
            f"Имя: {row['name']}\n"
            f"Фамилия: {row['surname']}\n"
            f"Город: {row['city']}\n"
            f"Стоимость: {row['cost']}\n"
            f"Часы: {row['hours']}\n"
            f"Мин. часы: {row['min_hours']}\n"
            f"Трансфер: {row['transfer']}\n"
            f"Instagram: {row['instagram']}\n"
            f"Портфолио: {row['portfolio']}\n"
            f"VK: {row['vk']}"
        )
        callback_data = f"select_{row['name']}_{row['surname']}" 
        if tuple(row.items()) in selected_personnel:
            button_text = f"✅ {button_text}"
        keyboard.append([InlineKeyboardButton(button_text, callback_data=callback_data)])
    
    keyboard.append([InlineKeyboardButton("Ввести название и дату", callback_data='enter_title_date')])
    reply_markup = InlineKeyboardMarkup(keyboard)
    if update.callback_query:
        await update.callback_query.message.reply_text('Выберите персонал:', reply_markup=reply_markup)
    else:
        await update.message.reply_text('Выберите персонал:', reply_markup=reply_markup)
    return "SELECTING_PERSONNEL"

async def select_personnel(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    query = update.callback_query
    await query.answer()
    person_id = query.data
    personnel_data = context.user_data.get('personnel_data', [])
    selected_personnel = context.user_data.get('selected_personnel', set())
    
    logging.info(f'Текущее состояние selected_personnel перед добавлением/удалением: {selected_personnel}')
    # Найдем выбранного пользователя
    selected_person = next((row for row in personnel_data if f"select_{row['name']}_{row['surname']}" == person_id), None)
    if selected_person:
        # Преобразуем dict в tuple of tuples для хранения в set
        selected_person_tuple = tuple(selected_person.items())
        if selected_person_tuple in selected_personnel:
            selected_personnel.remove(selected_person_tuple)
            logging.info(f'Пользователь удален из выбранных: {selected_person}')
        else:
            selected_personnel.add(selected_person_tuple)
            logging.info(f'Пользователь добавлен в выбранные: {selected_person}')
    logging.info(f'Текущее состояние selected_personnel после добавления/удаления: {selected_personnel}')
    context.user_data['selected_personnel'] = selected_personnel
    return await show_personnel_list(update, context)

async def choose_template(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    logging.info('Выбор шаблона')
    templates = get_template_names()
    if not templates:
        await update.callback_query.message.reply_text('Нет доступных шаблонов. Пожалуйста, загрузите шаблон.')
        return ConversationHandler.END
    keyboard = [[InlineKeyboardButton(template, callback_data=f"tpl_{idx}")] for idx, template in enumerate(templates)]
    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.callback_query.message.reply_text('Выберите шаблон:', reply_markup=reply_markup)
    return "CHOOSING_TEMPLATE"

async def select_template(update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    query = update.callback_query
    await query.answer()
    selected_template_idx = int(query.data.split('_')[1])
    templates = get_template_names()
    selected_template = templates[selected_template_idx]
    await query.edit_message_text(text=f"Вы выбрали шаблон: {selected_template}")
    
    # Убедитесь, что папка для обновлённых шаблонов существует
    output_directory = 'templates'
    os.makedirs(output_directory, exist_ok=True)
    
    # Получаем выбранный персонал и преобразуем его в список словарей
    selected_personnel = [dict(person) for person in context.user_data.get('selected_personnel', set())]
    
    # Получаем название и дату из контекста
    title = context.user_data.get('title', '')
    date = context.user_data.get('date', '')
    
    # Получаем название таблицы
    sheet_title = context.user_data.get('sheet_title', 'Без названия')
    
    fill_ppt_template(
        selected_personnel, 
        selected_template, 
        os.path.join(output_directory, f'updated_{selected_template}.pptx'), 
        title=title, 
        date=date,
        sheet_title=sheet_title 
    )
    logging.info(f'Данные добавлены в шаблон {selected_template}.pptx')
    
    # Конвертируем обновлённый PPTX в PDF
    pptx_template_path = os.path.join(output_directory, f'updated_{selected_template}.pptx')
    pdf_output_path = os.path.join(output_directory, f'updated_{selected_template}.pdf')
    convert_pptx_to_pdf(pptx_template_path, pdf_output_path)
    logging.info(f'Конвертация {pptx_template_path} в PDF завершена.')
    
    # Отправляем конвертированный PDF
    try:
        await query.message.reply_document(open(pdf_output_path, 'rb'))
    except TimedOut:
        logging.error('Время ожидания истекло при отправке документа. Повторная попытка...')
        await query.message.reply_document(open(pdf_output_path, 'rb'))
    return ConversationHandler.END


async def upload_template(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    await update.callback_query.message.reply_text('Пожалуйста, загрузите файл шаблона (формат .pptx).')

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    document = update.message.document
    if document.mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        logging.info('Получен документ формата .pptx')
        file = await document.get_file()
        file_path = os.path.join(ORIGINAL_TEMPLATES_DIR, document.file_name)
        await file.download_to_drive(file_path)
        logging.info(f'Документ загружен в {file_path}')
        
        # Отправляем подтверждение
        await update.message.reply_text(f'Шаблон {document.file_name} успешно загружен и сохранён.')

async def show_templates(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    templates = get_template_names()
    if templates:
        template_list = '\n'.join(templates)
        await update.callback_query.message.reply_text(f'Доступные шаблоны:\n{template_list}')
    else:
        await update.callback_query.message.reply_text('Нет доступных шаблонов.')

async def delete_template(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    templates = get_template_names()
    if templates:
        keyboard = [[InlineKeyboardButton(template, callback_data=f"del_{idx}")] for idx, template in enumerate(templates)]
        reply_markup = InlineKeyboardMarkup(keyboard)
        await update.callback_query.message.reply_text('Выберите шаблон для удаления:', reply_markup=reply_markup)
    else:
        await update.callback_query.message.reply_text('Нет доступных шаблонов для удаления.')

async def confirm_delete_template(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    query = update.callback_query
    await query.answer()
    selected_template_idx = int(query.data.split('_')[1])
    templates = get_template_names()
    selected_template = templates[selected_template_idx]
    file_path = os.path.join(ORIGINAL_TEMPLATES_DIR, f"{selected_template}.pptx")
    try:
        os.remove(file_path)
        await query.edit_message_text(text=f'Шаблон {selected_template} успешно удалён.')
        logging.info(f'Шаблон {selected_template} удалён.')
    except Exception as e:
        await query.edit_message_text(text=f'Ошибка при удалении шаблона: {str(e)}')
        logging.error(f'Ошибка при удалении шаблона {selected_template}: {e}')

async def button(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    query = update.callback_query
    await query.answer()

    if query.data == 'get_contractors':
        return await get_personnel_data(update, context)
    elif query.data == 'enter_title_date':
        return await ask_for_title(update, context)
    elif query.data == 'upload_template':
        await upload_template(update, context)
        return ConversationHandler.END
    elif query.data == 'show_templates':
        await show_templates(update, context)
        return ConversationHandler.END
    elif query.data == 'delete_template':
        await delete_template(update, context)
        return ConversationHandler.END
    elif query.data.startswith('del_'):
        await confirm_delete_template(update, context)
        return ConversationHandler.END
    elif query.data.startswith('select_'):
        return await select_personnel(update, context)
    elif query.data == 'choose_template':
        return await choose_template(update, context)
    elif query.data.startswith('tpl_'):
        return await select_template(update, context)
    else:
        return ConversationHandler.END

async def ask_for_title(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    logging.info('Запрос на ввод названия')
    if update.callback_query:
        await update.callback_query.message.reply_text('Пожалуйста, введите название:')
    else:
        await update.message.reply_text('Пожалуйста, введите название:')
    return "WAITING_FOR_TITLE"

async def receive_title(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    title = update.message.text
    logging.info(f'Получено название: {title}')
    context.user_data['title'] = title
    await update.message.reply_text('Теперь введите дату:')
    return "WAITING_FOR_DATE"

async def receive_date(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    date = update.message.text
    logging.info(f'Получена дата: {date}')
    context.user_data['date'] = date
    keyboard = [[InlineKeyboardButton("Выбрать шаблон", callback_data='choose_template')]]
    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.message.reply_text('Выберите действие:', reply_markup=reply_markup)
    return "CHOOSING_TEMPLATE"

async def get_personnel_data(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    spreadsheet_id = '1syI_hA_2n-4dygYYHZ5vt0_l2ys2uZXDlD3nONglWOU'
    range_name = 'Sheet1!A2:J'
    data = read_sheet(spreadsheet_id, range_name)
    logging.info(f"Извлеченные данные из Google Sheets: {data}")
    if not data:
        await update.callback_query.message.reply_text('Не удалось получить данные из таблицы.')
        return ConversationHandler.END
    
    # Получаем название таблицы
    sheet_title = get_sheet_title(spreadsheet_id)
    context.user_data['sheet_title'] = sheet_title  
    
    # Обработка данных для включения всех полей
    personnel_data = []
    for row in data:
        if len(row) < 10:
            row += ['N/A'] * (10 - len(row))
        personnel_data.append({
            'name': row[1] if len(row) > 1 and row[1] else 'Не указано', 
            'surname': row[1].split()[1] if len(row[1].split()) > 1 else 'N/A',
            'city': row[2],
            'cost': row[3],
            'hours': row[4],
            'min_hours': row[5],
            'transfer': row[6],
            'instagram': row[7],
            'portfolio': row[8],
            'vk': row[9],
            'photo': row[0] if len(row) > 0 else 'N/A' 
        })
    
    context.user_data['personnel_data'] = personnel_data
    logging.info(f'Извлеченные данные о персонале: {personnel_data}')
    # Отправляем список персонала и показываем кнопки выбора
    return await show_personnel_list(update, context)


def extract_username(url):
    if 'instagram.com' in url:
        parts = url.split('/')
        if len(parts) > 3:
            username = parts[3].split('?')[0]
            return username
    elif 'vk.com' in url:
        parts = url.split('/')
        if len(parts) > 3:
            username = parts[3].split('?')[0]
            return username
    return url 

def fill_ppt_template(selected_people, template_name, output_path='output.pptx', title='', date='', sheet_title=''):
    """
    Заполняет шаблон PPTX данными выбранных пользователей, добавляет название таблицы Google Sheets,
    название и дату, а затем сохраняет обновленную презентацию.

    :param selected_people: Список словарей с данными о пользователях.
    :param template_name: Название шаблона PPTX без расширения.
    :param output_path: Путь для сохранения обновленной презентации.
    :param title: Название презентации.
    :param date: Дата презентации.
    :param sheet_title: Название таблицы Google Sheets.
    """
    logging.info(f"Количество выбранных пользователей: {len(selected_people)}")
    logging.info(f'Данные для добавления в PPTX: {selected_people}')
    
    template_path = os.path.join(ORIGINAL_TEMPLATES_DIR, f"{template_name}.pptx")
    
    if not os.path.exists(template_path):
        logging.error(f"Шаблон {template_path} не найден.")
        raise FileNotFoundError(f"Шаблон {template_path} не найден.")
    
    presentation = Presentation(template_path)
    logging.info('Шаблон презентации загружен для заполнения')
    
    # Вставка названия и даты в первый слайд
    if len(presentation.slides) == 0:
        current_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    else:
        current_slide = presentation.slides[0]
    
    # Добавляем название
    title_shape = current_slide.shapes.title
    if title_shape:
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
    else:
        # Если заполнителя заголовка нет, добавляем текстовое поле
        left = Inches(2.05)
        top = Inches(11.34)
        width = presentation.slide_width - Inches(2)
        height = Inches(1)
        textbox = current_slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = 'Helvetica'
        run.font.size = Pt(20)
        run.font.bold = False
        run.font.color.rgb = RGBColor(255, 255, 255)

    # Добавляем дату
    left = Inches(3.15)
    top = Inches(12.15)
    width = Inches(2)
    height = Inches(0.5)
    date_shape = current_slide.shapes.add_textbox(left, top, width, height)
    date_frame = date_shape.text_frame
    p = date_frame.add_paragraph()
    run = p.add_run()
    run.text = date
    run.font.name = 'Helvetica'
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(255, 255, 255) 

    # Добавляем название таблицы Google Sheets
    if sheet_title:
        left = Inches(0.5)  # Левый отступ
        top = Inches(11.75)  # Верхний отступ
        width = Inches(6)    # Ширина текстового поля
        height = Inches(0.5) # Высота текстового поля
        sheet_title_shape = current_slide.shapes.add_textbox(left, top, width, height)
        sheet_title_frame = sheet_title_shape.text_frame
        p = sheet_title_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{sheet_title}"
        run.font.name = 'Helvetica'
        run.font.size = Pt(32)
        run.font.bold = False
        run.font.color.rgb = RGBColor(0, 0, 0) 
        logging.info('Добавлено текстовое поле для названия таблицы')

    logging.info('Используется первый слайд для размещения заголовка, даты и названия таблицы')
    
    # Параметры размещения
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    
    top_margin = Inches(1400 / 96) # Смещение вниз первого пользователя от него по логике остальные идут вниз
     
    left_margin_image = Inches(0.1)
    spacing_between_users = Inches(1) 

    image_width = Inches(2.6)
    image_height = Inches(3.6)
    textbox_width = Inches(2.9)
    textbox_height = Inches(2)

    current_top = top_margin 

    ICON_SHIFT_X = Inches(0.55)  # Смещение вправо
    ICON_SHIFT_Y = Inches(1.85)  # Смещение вниз
    

    for idx, person in enumerate(selected_people):
        logging.info(f'Обработка пользователя {idx + 1}/{len(selected_people)}: {person}')
        
        # Проверяем, поместится ли следующий пользователь на текущем слайде
        if current_top + max(image_height, textbox_height) > slide_height - Inches(0.5):
            # Добавляем новый слайд, если места недостаточно
            current_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            logging.info('Добавлен новый слайд для размещения дополнительных пользователей')
            current_top = Inches(1) 
        
        # Загрузка и вставка фото
        photo_url = person.get('photo', 'N/A')
        logging.info(f"Проверка URL фото для {person['name']} {person['surname']}: {photo_url}")
        if photo_url != 'N/A' and photo_url.startswith('http'):
            logging.info(f"Загрузка фото для {person['name']} {person['surname']} с {photo_url}")
            photo_path = download_photo(photo_url)
            if photo_path:
                logging.info(f"Вставка фото из {photo_path} в слайд")
                current_slide.shapes.add_picture(photo_path, left_margin_image, current_top, width=image_width, height=image_height)
                try:
                    os.remove(photo_path)
                    logging.info(f"Удалён временный файл изображения: {photo_path}")
                except OSError as e:
                    logging.error(f"Ошибка при удалении файла {photo_path}: {e}")
        
        textbox_left = left_margin_image + image_width + Inches(0.2)
        # Добавляем 10 пикселей в дюймах (~0.104 дюйма)
        textbox_top = current_top + Inches(22 / 96)
        
        textbox = current_slide.shapes.add_textbox(textbox_left, textbox_top, textbox_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.clear()
        
        text_frame.margin_left = Inches(0.52)  

        # Добавление информации о пользователе
        if person['name'] and person['name'] != 'Не указано':
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f"{person['name'].upper()}"
            run.font.name = 'Helvetica'
            run.font.size = Pt(18)
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0) 

        if person['city'] and person['city'] != 'N/A':
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f"/ {person['city']}"
            run.font.name = 'Helvetica'
            run.font.size = Pt(14)
            run.font.bold = False
            run.font.color.rgb = RGBColor(128, 128, 128)

        # Стоимость
        if person['cost'] and person['cost'] != 'N/A':
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f"{person['cost']} / час"
            run.font.name = 'Helvetica'
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0, 0, 0)
            p.line_spacing = Pt(25) 
            run.font.bold = False

        # Минимальные часы
        if person['min_hours'] and person['min_hours'] != 'N/A':
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f"Минимально от {person['min_hours']} часов"
            run.font.name = 'Helvetica'
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(128, 128, 128)
            run.font.bold = False

        # Трансфер
        if person['transfer'] and person['transfer'] != 'N/A':
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = f"+ {person['transfer']} / трансфер"
            run.font.name = 'Helvetica'
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0, 0, 0)
            run.font.bold = False

        # Социальные сети и портфолио с иконками
        social_media = [
            ('instagram', 'icons/i.png', RGBColor(0, 0, 0)),
            ('portfolio', 'icons/w.png', RGBColor(0, 0, 0)),
            ('vk', 'icons/v.png', RGBColor(0, 0, 0)),
        ]

        # Определяем начальное смещение для иконок относительно текста
        social_media_start_y = textbox_top
        social_media_line_height = Inches(0.33) 

        social_media_line_number = 0

        for key, icon_path, color in social_media:
            if person[key] and person[key] != 'N/A':
                if key == 'instagram':
                    username = extract_username(person[key])
                elif key == 'vk':
                    username = extract_username(person[key])
                else:
                    username = person[key]
                
                # Добавляем параграф для текста
                p = text_frame.add_paragraph()
                run = p.add_run()
                run.text = f"{username}" 
                run.font.name = 'Helvetica'
                run.font.size = Pt(20)
                run.font.color.rgb = color
                run.font.bold = False
                p.line_spacing = Pt(25) 

                # Определяем размер шрифта для расчета позиции иконки
                font_size_pt = run.font.size.pt if run.font.size else 14 
                
                # Вставка иконки
                if os.path.exists(icon_path):
                    icon_width = Inches(0.2)
                    icon_height = Inches(0.2)
                    # Позиционируем иконку слева от текста с добавлением смещений
                    icon_left = textbox_left - icon_width - Inches(0.1) + ICON_SHIFT_X 
                    icon_top = social_media_start_y + social_media_line_number * social_media_line_height + ICON_SHIFT_Y  # Смещение вниз
                    try:
                        current_slide.shapes.add_picture(icon_path, icon_left, icon_top, width=icon_width, height=icon_height)
                    except Exception as e:
                        logging.error(f"Ошибка при вставке иконки {icon_path}: {e}")
                else:
                    logging.warning(f"Иконка {icon_path} не найдена.")

                # Инкрементируем номер строки
                social_media_line_number += 1

        # Обновляем позицию для следующего пользователя
        current_top += max(image_height, textbox_height) + spacing_between_users

    # Сохранение презентации
    presentation.save(output_path)
    logging.info(f'Презентация сохранена как {output_path}')
        
def convert_pptx_to_pdf(input_pptx, output_pdf):
    """
    Конвертирует PPTX в PDF с помощью unoconv.
    """
    try:
        subprocess.run(['unoconv', '-f', 'pdf', '-o', output_pdf, input_pptx], check=True)
    except subprocess.CalledProcessError as e:
        logging.error(f"Ошибка при конвертации PPTX в PDF: {e}")
        raise

def main():
    TOKEN = "7995888417:AAH13VfCgV6XPYVOfCA-K-9ji4oLVZqj8GI" 
    application = ApplicationBuilder().token(TOKEN).build()

    conv_handler = ConversationHandler(
        entry_points=[CallbackQueryHandler(button)],
        states={
            "WAITING_FOR_TITLE": [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_title)],
            "WAITING_FOR_DATE": [MessageHandler(filters.TEXT & ~filters.COMMAND, receive_date)],
            "CHOOSING_TEMPLATE": [
                CallbackQueryHandler(select_template, pattern=r"^tpl_"),
                CallbackQueryHandler(choose_template, pattern='^choose_template$'),
            ],
            "SELECTING_PERSONNEL": [
                CallbackQueryHandler(select_personnel, pattern=r"^select_"),
                CallbackQueryHandler(button),
            ],
        },
        fallbacks=[],
    )

    application.add_handler(CommandHandler("start", start))
    application.add_handler(conv_handler)
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))

    logging.info('Бот запущен и готов к работе')
    application.run_polling()

if __name__ == '__main__':
    main()